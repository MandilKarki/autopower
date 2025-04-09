from pptx import Presentation
from pptx.util import Inches

# Create a presentation
prs = Presentation()

# Add title slide
slide = prs.slides.add_slide(prs.slide_layouts[0])
title = slide.shapes.title
subtitle = slide.placeholders[1]

title.text = "Presentation Title"
subtitle.text = "Subtitle"

# Add content slide
slide = prs.slides.add_slide(prs.slide_layouts[1])
title = slide.shapes.title
content = slide.placeholders[1]

title.text = "Slide Title"
content.text = "Content goes here"

# Save the template
prs.save("templates/template.pptx")
