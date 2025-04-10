import os
import json
import csv
import numpy as np
import pandas as pd
import matplotlib.pyplot as plt
import seaborn as sns
import pdfplumber
import chromadb
import ollama
import matplotlib
import re
import random
from collections import Counter
from datetime import datetime
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from sklearn.feature_extraction.text import TfidfVectorizer
from sklearn.cluster import KMeans
from sklearn.decomposition import PCA
matplotlib.use('Agg')  # Use non-interactive backend

class IntelligentReportGenerator:
    """AI-Powered Intelligent Report Generator with PowerPoint and PDF capabilities"""
    
    def __init__(self, report_title="Intelligent Report", model_name=None):
        """Initialize the Intelligent Report Generator"""
        self.report_title = report_title
        self.sections = {}
        self.visualizations = []
        self.referenced_documents = []
        self.dataset = None
        self.available_models = self._get_available_models()
        
        # Set default model if none provided or if provided model isn't available
        if model_name is None or model_name not in self.available_models:
            self.model_name = self.available_models[0] if self.available_models else "llama2"
            if model_name is not None and model_name not in self.available_models:
                print(f"Warning: Model {model_name} not available. Using {self.model_name} instead.")
        else:
            self.model_name = model_name
        
        print(f"Using AI model: {self.model_name}")
        
        # Initialize ChromaDB for document embeddings
        self.chroma_client = chromadb.Client()
        try:
            self.collection = self.chroma_client.create_collection("report_documents")
        except:
            # Collection might already exist
            self.collection = self.chroma_client.get_collection("report_documents")
        
        # Create output directories if they don't exist
        os.makedirs("output", exist_ok=True)
        os.makedirs("temp", exist_ok=True)
        
    def _get_available_models(self):
        """Get list of available Ollama models"""
        try:
            response = ollama.list()
            return [model['name'] for model in response['models']]
        except Exception as e:
            print(f"Warning: Could not retrieve model list from Ollama: {e}")
            return []
    
    def add_user_section(self, section_title, content):
        """Add a user-provided section to the report"""
        self.sections[section_title] = {"type": "user", "content": content}
        return self
    
    def generate_ai_section(self, section_title, prompt, detailed=False):
        """Generate an AI-written section based on a prompt"""
        # Direct call to Ollama API
        word_count = "400-600" if detailed else "200-300"
        full_prompt = f'''Write a professional, informative section about {prompt}.
        The content should be factual, well-structured, and approximately {word_count} words.
        Use a formal tone suitable for a business or academic report.
        Include relevant statistics or examples where appropriate.
        Organize the content with clear paragraph breaks.'''
        
        response = ollama.chat(model=self.model_name, messages=[
            {
                "role": "user",
                "content": full_prompt
            }
        ])
        
        result = response['message']['content']
        self.sections[section_title] = {"type": "ai", "content": result}
        return self
    
    def suggest_topics(self, base_topic="business analytics", count=5):
        """Suggest relevant topics for report sections based on a base topic"""
        prompt = f'''Suggest {count} specific, focused topics related to "{base_topic}" that would make good sections in a professional report.
        For each topic, provide: 
        1. A concise section title (3-6 words)
        2. A brief description of what this section should cover (1-2 sentences)
        
        Format your response as a JSON list with objects containing 'title' and 'description' fields.
        Example: [{{'title': 'Market Segmentation Analysis', 'description': 'Examination of key market segments and targeting strategies.'}}]'''
        
        response = ollama.chat(model=self.model_name, messages=[
            {
                "role": "user",
                "content": prompt
            }
        ])
        
        # Extract JSON from response
        result = response['message']['content']
        
        # Try to extract JSON if it's embedded in markdown code blocks
        json_matches = re.findall(r'```(?:json)?\s*(.+?)\s*```', result, re.DOTALL)
        if json_matches:
            result = json_matches[0]
            
        try:
            topics = json.loads(result)
            return topics
        except json.JSONDecodeError:
            # If JSON parsing fails, extract manually
            print("Warning: Could not parse JSON response. Returning raw text.")
            return {"error": "Could not parse topics", "raw_response": result}
    
    def load_data(self, file_path):
        """Load data from CSV, Excel, or JSON file for analysis and visualization"""
        if not os.path.exists(file_path):
            print(f"Error: File {file_path} does not exist.")
            return self
            
        file_ext = os.path.splitext(file_path)[1].lower()
        
        try:
            if file_ext == '.csv':
                self.dataset = pd.read_csv(file_path)
            elif file_ext in ['.xls', '.xlsx']:
                self.dataset = pd.read_excel(file_path)
            elif file_ext == '.json':
                self.dataset = pd.read_json(file_path)
            else:
                print(f"Error: Unsupported file format {file_ext}")
                return self
                
            print(f"Data loaded successfully: {self.dataset.shape[0]} rows, {self.dataset.shape[1]} columns")
            
            # Auto-generate reports about the data
            self._analyze_dataset()
            
            return self
        except Exception as e:
            print(f"Error loading data: {e}")
            return self
    
    def _analyze_dataset(self):
        """Analyze loaded dataset and store insights"""
        if self.dataset is None or self.dataset.empty:
            return
            
        # Store dataset summary
        self.data_summary = {
            "rows": self.dataset.shape[0],
            "columns": self.dataset.shape[1],
            "column_types": {col: str(dtype) for col, dtype in self.dataset.dtypes.items()},
            "missing_values": self.dataset.isnull().sum().to_dict(),
            "numeric_columns": self.dataset.select_dtypes(include=[np.number]).columns.tolist(),
            "categorical_columns": self.dataset.select_dtypes(include=['object']).columns.tolist()
        }
        
        # Auto-suggest visualizations based on data types
        self._suggest_visualizations()
    
    def _suggest_visualizations(self):
        """Automatically suggest visualizations based on dataset structure"""
        if self.dataset is None or self.dataset.empty:
            return []
            
        suggestions = []
        
        numeric_cols = self.data_summary["numeric_columns"]
        categorical_cols = self.data_summary["categorical_columns"]
        
        # Suggestion 1: If we have numeric columns, suggest distributions
        if numeric_cols:
            for col in numeric_cols[:3]:  # Limit to first 3 to avoid too many suggestions
                suggestions.append({
                    "title": f"Distribution of {col}",
                    "type": "histogram",
                    "columns": [col],
                    "description": f"Histogram showing the distribution of {col} values."
                })
                
        # Suggestion 2: If we have categorical columns, suggest count plots
        if categorical_cols:
            for col in categorical_cols[:3]:  # Limit to first 3
                suggestions.append({
                    "title": f"Count of {col} Categories",
                    "type": "bar",
                    "columns": [col],
                    "description": f"Bar chart showing the count of each {col} category."
                })
                
        # Suggestion 3: If we have both numeric and categorical, suggest relationship
        if numeric_cols and categorical_cols:
            suggestions.append({
                "title": f"Relationship: {categorical_cols[0]} vs {numeric_cols[0]}",
                "type": "boxplot",
                "columns": [categorical_cols[0], numeric_cols[0]],
                "description": f"Box plot showing {numeric_cols[0]} distribution across {categorical_cols[0]} categories."
            })
            
        # Suggestion 4: If we have multiple numeric columns, suggest correlation
        if len(numeric_cols) >= 2:
            suggestions.append({
                "title": "Correlation Matrix",
                "type": "heatmap",
                "columns": numeric_cols[:5],  # Limit to 5 columns for readability
                "description": "Heatmap showing correlations between numeric variables."
            })
            
            suggestions.append({
                "title": f"Scatter Plot: {numeric_cols[0]} vs {numeric_cols[1]}",
                "type": "scatter",
                "columns": [numeric_cols[0], numeric_cols[1]],
                "description": f"Scatter plot showing relationship between {numeric_cols[0]} and {numeric_cols[1]}."
            })
            
        return suggestions
        
    def ingest_document(self, file_path):
        """Ingest a document for retrieval processing"""
        if not os.path.exists(file_path):
            print(f"Error: File {file_path} does not exist.")
            return self
        
        # Store reference to the document
        self.referenced_documents.append(file_path)
        
        # Extract text from PDF or text file
        if file_path.lower().endswith('.pdf'):
            text = self._extract_text_from_pdf(file_path)
        else:
            # For text files
            with open(file_path, 'r', encoding='utf-8') as f:
                text = f.read()
        
        # Simple chunking approach
        chunk_size = 1000
        overlap = 200
        chunks = []
        
        # Split the text into chunks with overlap
        for i in range(0, len(text), chunk_size - overlap):
            chunk = text[i:i + chunk_size]
            if chunk.strip():
                chunks.append(chunk)
        
        # Get embeddings from Ollama and store in ChromaDB
        for i, chunk in enumerate(chunks):
            # Get embedding from Ollama
            embedding_response = ollama.embeddings(
                model=self.model_name,
                prompt=chunk
            )
            
            # Store in ChromaDB
            self.collection.add(
                ids=[f"{file_path}_{i}"],
                documents=[chunk],
                metadatas=[{"source": file_path, "chunk": i}]
            )
        
        # Extract key topics from document for suggesting report sections
        topics = self._extract_topics_from_text(text)
        print(f"Extracted key topics: {', '.join(topics[:5])}")
        
        return self
        
    def _extract_topics_from_text(self, text, num_topics=5):
        """Extract key topics from a document using TF-IDF"""
        # Clean and tokenize text
        sentences = [s.strip() for s in re.split(r'[.!?]', text) if s.strip()]
        
        if len(sentences) < 5:
            # Not enough text for meaningful extraction
            return []
            
        try:
            # Create TF-IDF vectors
            vectorizer = TfidfVectorizer(max_df=0.7, min_df=2, stop_words='english')
            tfidf_matrix = vectorizer.fit_transform(sentences)
            
            # Get top terms
            feature_names = vectorizer.get_feature_names_out()
            tfidf_scores = tfidf_matrix.sum(axis=0).A1
            top_indices = tfidf_scores.argsort()[-num_topics*2:][::-1]  # Get more than needed to filter
            top_terms = [feature_names[i] for i in top_indices]
            
            # Filter to more meaningful terms (longer words tend to be more meaningful)
            top_terms = [term for term in top_terms if len(term) > 3][:num_topics]
            
            return top_terms
        except Exception as e:
            print(f"Error extracting topics: {e}")
            return []
    
    def rag_query(self, section_title, query):
        """Query ingested documents and add the results as a section"""
        if not self.referenced_documents:
            self.sections[section_title] = {
                "type": "rag",
                "content": "No documents have been ingested for processing."
            }
            return self
        
        # Get query embedding from Ollama
        query_embedding = ollama.embeddings(
            model=self.model_name,
            prompt=query
        )
        
        # Query ChromaDB for similar documents
        results = self.collection.query(
            query_embeddings=query_embedding['embedding'],
            n_results=3
        )
        
        if not results['documents'][0]:
            self.sections[section_title] = {
                "type": "rag",
                "content": f"No relevant information found for: {query}"
            }
            return self
        
        # Format retrieved information
        context = "\n\n".join(results['documents'][0])
        sources = "\n".join([f"• {meta['source']}" for meta in results['metadatas'][0]])
        
        # Generate the section content with Ollama
        prompt = f'''Based on the following information:

{context}


        Answer the query: {query}

Provide a well-structured response that synthesizes 
        the information in a way that directly addresses the query.'''
        
        response = ollama.chat(model=self.model_name, messages=[
            {
                "role": "user",
                "content": prompt
            }
        ])
        
        result = response['message']['content']
        
        # Create section with citations
        content = f"{result}\n\nSources:\n{sources}"
        self.sections[section_title] = {"type": "rag", "content": content}
        
        return self
    
    def add_visualization(self, title, data, chart_type="bar", labels=None, **kwargs):
        """Add a data visualization to the report using raw data"""
        fig, ax = plt.subplots(figsize=(10, 6))
        
        # Extract labels for axes - remove from kwargs to avoid passing to plot functions
        xlabel = kwargs.pop('xlabel', None)
        ylabel = kwargs.pop('ylabel', None)
        description = kwargs.pop('description', '')
        
        # Create the specified chart type
        if chart_type.lower() == "bar":
            if labels is None:
                labels = list(range(len(data)))
            ax.bar(labels, data, **kwargs)
            plt.xticks(rotation=45, ha='right')
        elif chart_type.lower() == "line":
            if labels is None:
                labels = list(range(len(data)))
            ax.plot(labels, data, marker='o', **kwargs)
            plt.xticks(rotation=45, ha='right')
        elif chart_type.lower() == "pie":
            if labels is None:
                labels = [f"Segment {i+1}" for i in range(len(data))]
            ax.pie(data, labels=labels, autopct='%1.1f%%', shadow=True, **kwargs)
        elif chart_type.lower() == "scatter":
            # For scatter, data should be a list of (x,y) tuples
            x, y = zip(*data)
            ax.scatter(x, y, **kwargs)
        elif chart_type.lower() == "histogram":
            bins = kwargs.pop('bins', 10)
            ax.hist(data, bins=bins, **kwargs)
        
        # Add labels if provided
        if xlabel:
            ax.set_xlabel(xlabel)
        if ylabel:
            ax.set_ylabel(ylabel)
            
        # Add title and style
        ax.set_title(title)
        plt.tight_layout()
        
        # Save the figure
        filename = f"temp/viz_{len(self.visualizations)}.png"
        plt.savefig(filename, dpi=300)
        plt.close()
        
        # Store visualization metadata
        self.visualizations.append({
            "title": title,
            "filename": filename,
            "chart_type": chart_type,
            "description": description
        })
        
        return self
        
    def create_visualization_from_dataframe(self, title, chart_type, columns=None, **kwargs):
        """Create visualization from loaded dataframe data"""
        if self.dataset is None or self.dataset.empty:
            print("Error: No dataset loaded. Use load_data() first.")
            return self
            
        df = self.dataset
        
        # Create figure
        fig, ax = plt.subplots(figsize=(10, 6))
        
        # Apply style
        style = kwargs.get('style', 'seaborn-v0_8-whitegrid')
        try:
            plt.style.use(style)
        except:
            pass  # Use default style if specified style not available
            
        # Create visualization based on chart type
        if chart_type.lower() == 'bar':
            if len(columns) == 1:
                # Simple count plot
                col = columns[0]
                counts = df[col].value_counts()
                counts.plot(kind='bar', ax=ax, **{k:v for k,v in kwargs.items() if k not in ['style', 'description']})
                ax.set_xlabel(col)
                ax.set_ylabel('Count')
            elif len(columns) == 2:
                # Grouped bar chart
                df.groupby(columns[0])[columns[1]].mean().plot(kind='bar', ax=ax)
                ax.set_ylabel(f'Mean {columns[1]}')
                
        elif chart_type.lower() == 'line':
            if 'x' in kwargs and 'y' in kwargs:
                df.plot(x=kwargs['x'], y=kwargs['y'], kind='line', ax=ax, marker='o')
            elif len(columns) >= 1:
                # Plot time series or sequential data
                if len(columns) == 1:
                    df[columns[0]].plot(kind='line', ax=ax, marker='o')
                else:
                    df[columns].plot(kind='line', ax=ax, marker='o')
                    
        elif chart_type.lower() == 'scatter':
            if len(columns) >= 2:
                x, y = columns[0], columns[1]
                color_col = columns[2] if len(columns) > 2 else None
                
                if color_col:
                    # Use third column for color coding
                    scatter = ax.scatter(df[x], df[y], c=df[color_col], cmap='viridis', alpha=0.7)
                    plt.colorbar(scatter, ax=ax, label=color_col)
                else:
                    ax.scatter(df[x], df[y], alpha=0.7)
                    
                ax.set_xlabel(x)
                ax.set_ylabel(y)
                
        elif chart_type.lower() == 'histogram':
            if len(columns) >= 1:
                col = columns[0]
                ax.hist(df[col], bins=kwargs.get('bins', 10), alpha=0.7)
                ax.set_xlabel(col)
                ax.set_ylabel('Frequency')
                
        elif chart_type.lower() == 'boxplot':
            if len(columns) == 1:
                # Single boxplot
                df.boxplot(column=columns[0], ax=ax)
            elif len(columns) == 2:
                # Grouped boxplot
                df.boxplot(column=columns[1], by=columns[0], ax=ax)
                plt.suptitle('')  # Remove default title
                ax.set_title(title)
                
        elif chart_type.lower() == 'heatmap':
            if len(columns) >= 2:
                # Correlation heatmap
                selected_df = df[columns] if columns else df.select_dtypes(include=[np.number])
                corr = selected_df.corr()
                sns.heatmap(corr, annot=True, cmap='coolwarm', linewidths=0.5, ax=ax)
                
        elif chart_type.lower() == 'pie':
            if len(columns) == 1:
                col = columns[0]
                counts = df[col].value_counts()
                # Limit pie chart to top categories if there are many
                if len(counts) > 8:
                    counts_limited = counts.nlargest(7)
                    counts_limited['Other'] = counts[7:].sum()
                    counts_limited.plot(kind='pie', autopct='%1.1f%%', ax=ax)
                else:
                    counts.plot(kind='pie', autopct='%1.1f%%', ax=ax)
                ax.set_ylabel('')  # Remove ylabel
        
        # Add title
        ax.set_title(title)
        plt.tight_layout()
        
        # Save the figure
        filename = f"temp/viz_{len(self.visualizations)}.png"
        plt.savefig(filename, dpi=300, bbox_inches='tight')
        plt.close()
        
        # Store visualization metadata
        self.visualizations.append({
            "title": title,
            "filename": filename,
            "chart_type": chart_type,
            "columns": columns,
            "description": kwargs.get('description', '')
        })
        
        return self
    
    def auto_visualize(self, count=3):
        """Automatically generate the most meaningful visualizations from the dataset"""
        if self.dataset is None or self.dataset.empty:
            print("Error: No dataset loaded. Use load_data() first.")
            return self
            
        # Get suggested visualizations
        suggestions = self._suggest_visualizations()
        if not suggestions:
            return self
            
        # Create the top N suggested visualizations
        for i, suggestion in enumerate(suggestions[:count]):
            self.create_visualization_from_dataframe(
                title=suggestion['title'],
                chart_type=suggestion['type'],
                columns=suggestion['columns'],
                description=suggestion['description']
            )
            
        return self
    
    def _extract_text_from_pdf(self, pdf_path):
        """Extract text from a PDF file"""
        text = ""
        with pdfplumber.open(pdf_path) as pdf:
            for page in pdf.pages:
                text += page.extract_text() + "\n"
        return text
    
    def _format_section(self, section_name, section_data):
        """Format a section for the PDF"""
        result = f"# {section_name}\n\n"
        
        if section_data["type"] == "user" or section_data["type"] == "ai" or section_data["type"] == "rag":
            result += section_data["content"]
        
        return result
    
    def generate_report(self, output_path="output/generated_report.md"):
        """Generate the final report in Markdown format"""
        # Compile the report content
        content = [f"# {self.report_title}\n"]
        
        # Add a table of contents
        content.append("## Table of Contents\n")
        for i, section_name in enumerate(self.sections.keys(), 1):
            content.append(f"{i}. {section_name}\n")
        content.append("\n")
        
        # Add each section
        for section_name, section_data in self.sections.items():
            content.append(self._format_section(section_name, section_data))
            content.append("\n\n")
        
        # Add visualizations
        if self.visualizations:
            content.append("## Visualizations\n")
            for viz in self.visualizations:
                content.append(f"### {viz['title']}\n")
                content.append(f"![{viz['title']}]({viz['filename']})\n\n")
        
        # Write to file
        with open(output_path, "w", encoding="utf-8") as f:
            f.write("\n".join(content))
        
        print(f"Report generated successfully: {output_path}")
        return output_path

    def generate_presentation(self, output_path="output/presentation.pptx"):
        """Generate a PowerPoint presentation based on the report data"""
        # Create a new presentation
        prs = Presentation()
        
        # Add title slide
        title_slide_layout = prs.slide_layouts[0]  # Title slide layout
        slide = prs.slides.add_slide(title_slide_layout)
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        
        title.text = self.report_title
        subtitle.text = "Generated on " + os.path.basename(output_path).split('.')[0]
        
        # Add a slide for each section
        for section_name, section_data in self.sections.items():
            # Create section slide with heading and content
            content_slide_layout = prs.slide_layouts[1]  # Title and content layout
            slide = prs.slides.add_slide(content_slide_layout)
            
            # Set the title
            title = slide.shapes.title
            title.text = section_name
            
            # Add content
            content = slide.placeholders[1]
            
            # Format the content based on the section type
            if section_data["type"] == "user":
                content.text = section_data["content"]
            elif section_data["type"] == "ai":
                # Limit AI-generated content to fit in a slide
                text = section_data["content"]
                # Split into paragraphs and limit to first 3-4 paragraphs or ~150 words
                paragraphs = text.split('\n\n')
                content.text = '\n'.join(paragraphs[:min(3, len(paragraphs))])
                
                # If content is too long, add a "continued..." slide
                if len(paragraphs) > 3:
                    slide = prs.slides.add_slide(content_slide_layout)
                    title = slide.shapes.title
                    title.text = f"{section_name} (Continued)"
                    content = slide.placeholders[1]
                    content.text = '\n'.join(paragraphs[3:min(6, len(paragraphs))])
            elif section_data["type"] == "rag":
                # For RAG content, include the main content and sources
                text_parts = section_data["content"].split("Sources:")
                if len(text_parts) > 1:
                    # Main content
                    main_text = text_parts[0].strip()
                    # Sources
                    sources = text_parts[1].strip()
                    
                    # Add main content
                    paragraphs = main_text.split('\n\n')
                    content.text = '\n'.join(paragraphs[:min(3, len(paragraphs))])
                    
                    # Add sources in a bullet list
                    txBox = slide.shapes.add_textbox(Inches(0.5), Inches(5), Inches(9), Inches(1.5))
                    tf = txBox.text_frame
                    tf.text = "Sources:"
                    
                    # Add sources as bullet points
                    for source_line in sources.split('\n'):
                        p = tf.add_paragraph()
                        p.text = source_line.strip()
                        p.level = 1  # Set bullet level
                else:
                    # If no sources, just add the content
                    content.text = section_data["content"]
        
        # Add visualization slides
        for viz in self.visualizations:
            # Create slide with just a title
            title_only_slide_layout = prs.slide_layouts[5]  # Title only layout
            slide = prs.slides.add_slide(title_only_slide_layout)
            
            # Set the title
            title = slide.shapes.title
            title.text = viz['title']
            
            # Add the visualization image
            if os.path.exists(viz['filename']):
                left = Inches(1.5)
                top = Inches(2)
                width = Inches(7)  # Adjust based on image aspect ratio
                slide.shapes.add_picture(viz['filename'], left, top, width=width)
        
        # Save the presentation
        os.makedirs(os.path.dirname(output_path), exist_ok=True)
        prs.save(output_path)
        
        print(f"PowerPoint presentation generated successfully: {output_path}")
        return output_path

# Example usage
if __name__ == "__main__":
    # Create a new intelligent report generator (will use best available model)
    report = IntelligentReportGenerator("Market Analysis Report 2024")
    
    # Option 1: Get topic suggestions based on a general area
    print("\nSuggested report topics:\n======================")
    suggested_topics = report.suggest_topics("market analysis", count=3)
    for i, topic in enumerate(suggested_topics, 1):
        print(f"{i}. {topic['title']}: {topic['description']}")
    
    # Option 2: Create a sample dataset if you don't have one
    # In a real scenario, you'd use report.load_data("your_data.csv") instead
    import pandas as pd
    sample_data = {
        "Sector": ["Tech", "Healthcare", "Finance", "Retail", "Energy"] * 5,
        "Revenue": [random.randint(50, 200) for _ in range(25)],
        "Growth": [random.uniform(-5, 25) for _ in range(25)],
        "Employees": [random.randint(100, 5000) for _ in range(25)],
        "Satisfaction": [random.uniform(3.0, 5.0) for _ in range(25)]
    }
    df = pd.DataFrame(sample_data)
    df.to_csv("temp/sample_data.csv", index=False)
    
    # Load the sample data
    report.load_data("temp/sample_data.csv")
    
    # Automatically generate visualizations based on the data
    report.auto_visualize(count=3)
    
    # Add user-provided content
    report.add_user_section(
        "Executive Summary",
        "This report provides a comprehensive analysis of market trends across key sectors in 2024. "
        "The findings presented here are based on data collected from multiple sources and analyzed using "
        "advanced statistical methods and AI techniques."
    )
    
    # Generate AI content (using the detailed option for more comprehensive content)
    report.generate_ai_section(
        "Industry Overview",
        "Current state of the technology industry in 2024, including major trends, disruptions, and growth sectors.",
        detailed=True
    )
    
    # Add custom visualizations if needed
    market_data = [42, 35, 63, 27, 45]
    market_labels = ["Tech", "Healthcare", "Finance", "Retail", "Energy"]
    report.add_visualization(
        "Market Growth Forecast",
        market_data,
        "bar",
        market_labels,
        color="skyblue",
        xlabel="Industry Sector",
        ylabel="Projected Growth (%)",
        description="Projected growth percentages across major industry sectors for 2024-2025."
    )
    
    # Add more customized visualizations using dataframe data
    report.create_visualization_from_dataframe(
        title="Sector Performance Comparison",
        chart_type="boxplot",
        columns=["Sector", "Growth"],
        description="Distribution of growth rates across different sectors showing median, quartiles, and outliers."
    )
    
    # Generate the markdown report
    report.generate_report("output/market_analysis_report.md")
    
    # Generate the PowerPoint presentation
    report.generate_presentation("output/market_analysis_presentation.pptx")
    
    print("\nGenerated files:\n===============")
    print("1. Report (Markdown): output/market_analysis_report.md")
    print("2. Presentation (PowerPoint): output/market_analysis_presentation.pptx")
    print("\nNote: Visualizations are stored in the 'temp' directory.")
