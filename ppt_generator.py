import os
import json
import re
import uuid
from pathlib import Path
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from io import BytesIO
import requests
import pandas as pd
import numpy as np
import matplotlib.pyplot as plt
from PIL import Image
from PyPDF2 import PdfReader
import tempfile

# For RAG implementation
import chromadb
from chromadb.config import Settings
from chromadb.utils import embedding_functions
import openai
from sentence_transformers import SentenceTransformer

class RagDocumentProcessor:
    def __init__(self, collection_name=None):
        """
        Initialize the RAG document processor with embeddings and vector store
        
        Args:
            collection_name (str, optional): Name for the ChromaDB collection
        """
        # Use Ollama for embeddings
        import requests
        
        # Create a custom embedding function class for Ollama
        class OllamaEmbedding:
            def __init__(self, model="nomic-embed-text", base_url="http://localhost:11434"):
                self.model = model
                self.base_url = base_url
                self.api_endpoint = f"{base_url}/api/embeddings"
                
                # Test connection
                try:
                    r = requests.get(f"{base_url}/api/tags")
                    if r.status_code == 200:
                        print(f"Successfully connected to Ollama at {base_url}")
                    else:
                        print(f"Warning: Ollama is available but returned status {r.status_code}")
                except Exception as e:
                    print(f"Warning: Could not connect to Ollama at {base_url}: {str(e)}")
                    
            def __call__(self, input):
                # Handle both single strings and lists of strings
                if isinstance(input, str):
                    texts = [input]
                else:
                    texts = input
                    
                all_embeddings = []
                
                for text in texts:
                    try:
                        response = requests.post(
                            self.api_endpoint,
                            json={"model": self.model, "prompt": text}
                        )
                        
                        if response.status_code == 200:
                            embedding = response.json().get("embedding")
                            if embedding:
                                all_embeddings.append(embedding)
                            else:
                                # If no embedding returned, use a fallback approach
                                print(f"Warning: Ollama returned no embedding for text: {text[:50]}...")
                                all_embeddings.append([0.0] * 384)  # Use zero vector as fallback
                        else:
                            print(f"Warning: Ollama embedding request failed with status {response.status_code}")
                            all_embeddings.append([0.0] * 384)  # Use zero vector as fallback
                    except Exception as e:
                        print(f"Error getting embedding from Ollama: {str(e)}")
                        all_embeddings.append([0.0] * 384)  # Use zero vector as fallback
                        
                # For single inputs, return just the embedding
                if isinstance(input, str):
                    return all_embeddings[0]
                    
                return all_embeddings
        
        # Use Ollama for embeddings
        self.embedding_function = OllamaEmbedding(model="nomic-embed-text")
        self.using_ollama = True
            
        # Initialize ChromaDB client and collection using the latest configuration approach
        self.client = chromadb.PersistentClient(path="chroma_db")
        
        # Create or get collection
        if not collection_name:
            collection_name = f"docs_{uuid.uuid4().hex[:8]}"
        
        try:
            self.collection = self.client.get_or_create_collection(
                name=collection_name,
                embedding_function=self.embedding_function
            )
        except:
            # If there's an issue with the existing collection, create a new one
            collection_name = f"docs_{uuid.uuid4().hex[:8]}"
            self.collection = self.client.get_or_create_collection(
                name=collection_name, 
                embedding_function=self.embedding_function
            )
        
        self.added_docs = 0
    
    def process_pdf(self, pdf_path, chunk_size=1000, chunk_overlap=200):
        """
        Process a PDF document, split into chunks, and add to vector database
        
        Args:
            pdf_path (str): Path to the PDF file
            chunk_size (int): Size of text chunks for embedding
            chunk_overlap (int): Overlap between chunks
            
        Returns:
            int: Number of chunks added to the database
        """
        # Extract text from PDF
        with open(pdf_path, 'rb') as file:
            reader = PdfReader(file)
            text = ''
            for page in reader.pages:
                text += page.extract_text() + '\n\n'
        
        # Split text into chunks with overlap
        chunks = []
        for i in range(0, len(text), chunk_size - chunk_overlap):
            chunk = text[i:i + chunk_size]
            if len(chunk) >= 100:  # Only add chunks with meaningful content
                chunks.append(chunk)
        
        # Add chunks to vector database
        if chunks:
            ids = [f"doc_{self.added_docs + i}" for i in range(len(chunks))]
            metadatas = [{
                "source": os.path.basename(pdf_path),
                "chunk": i
            } for i in range(len(chunks))]
            
            self.collection.add(
                documents=chunks,
                ids=ids,
                metadatas=metadatas
            )
            
            self.added_docs += len(chunks)
            return len(chunks)
        
        return 0
    
    def query_documents(self, query, n_results=5):
        """
        Query the vector database with a question
        
        Args:
            query (str): Query text
            n_results (int): Number of results to return
            
        Returns:
            list: List of relevant document chunks
        """
        if self.added_docs == 0:
            return []
            
        results = self.collection.query(
            query_texts=[query],
            n_results=min(n_results, self.added_docs)
        )
        
        if not results['documents'][0]:
            return []
            
        return [
            {
                "content": doc,
                "metadata": meta
            } for doc, meta in zip(results['documents'][0], results['metadatas'][0])
        ]
    
    def generate_content_with_rag(self, query, system_prompt=None, model="llama3.2:latest", ollama_base_url="http://localhost:11434"):
        """
        Generate content using RAG by querying context and then an LLM
        
        Args:
            query (str): The query or instruction for generating content
            system_prompt (str, optional): System prompt for the LLM
            model (str): Model to use for generation (for Ollama: "llama3", "mistral", etc.)
            ollama_base_url (str): Base URL for Ollama API if using Ollama models
            
        Returns:
            str: Generated content
        """
        # Get relevant documents
        docs = self.query_documents(query)
        if not docs:
            context = "No relevant information found in the provided documents."
        else:
            context = "\n\n".join([f"DOCUMENT {i+1}:\n{doc['content']}" for i, doc in enumerate(docs)])
        
        # Default system prompt if none provided
        if not system_prompt:
            system_prompt = """You are a helpful assistant that generates presentation content.
            Use the provided context to create informative and engaging content for a PowerPoint slide.
            Make sure the content is concise, well-structured, and appropriate for a presentation format."""
            
        prompt = f"Here is relevant context from documents:\n{context}\n\nBased on this context, please generate content for: {query}"
        
        try:
            # First try using Ollama with the provided model (default: llama3)
            try:
                import requests
                
                # Check if ollama_base_url is accessible
                try:
                    # Simple ping to check if Ollama is running
                    ping_response = requests.get(f"{ollama_base_url}/api/tags", timeout=2)
                    if ping_response.status_code == 200:
                        # Prepare the request for Ollama
                        ollama_request = {
                            "model": model,
                            "prompt": prompt,
                            "system": system_prompt,
                            "stream": False
                        }
                        
                        # Make the request to Ollama
                        response = requests.post(
                            f"{ollama_base_url}/api/generate",
                            json=ollama_request,
                            timeout=30  # Longer timeout for generation
                        )
                        
                        if response.status_code == 200:
                            return response.json().get("response", "Ollama returned no content")
                except Exception as ollama_error:
                    print(f"Ollama not available or error: {str(ollama_error)}")
            except ImportError:
                print("Requests library not available for Ollama API calls")
            
            # If Ollama is not available, just return the context
            # No OpenAI fallback - using 100% local models
            return f"Content for {query}:\n\n{context[:2000]}... (content truncated)"
        except Exception as e:
            return f"Error generating content: {str(e)}\n\nHere's the context found:\n{context[:1000]}... (content truncated)"

class PowerPointGenerator:
    def __init__(self, template_path=None, model="llama3.2:latest"):
        """
        Initialize the PowerPoint generator with an optional template
        
        Args:
            template_path (str, optional): Path to the PowerPoint template file
            model (str, optional): Ollama model to use for content generation (default: llama3.2:latest)
        """
        if template_path and os.path.exists(template_path):
            if template_path.lower().endswith('.pptx'):
                self.presentation = Presentation(template_path)
            else:
                # Not a PowerPoint file, create new presentation
                self.presentation = Presentation()
        else:
            self.presentation = Presentation()
        
        # Initialize RAG processor for document context with the specified model
        self.rag_processor = RagDocumentProcessor()
        self.model = model  # Store the model name for later use
        
        self.slide_structures = self._analyze_template() if template_path and template_path.lower().endswith('.pptx') else []
    
    def _analyze_template(self):
        """
        Analyze the template structure to understand placeholders and layouts
        
        Returns:
            list: Structured information about slides and their placeholders
        """
        slide_structures = []
        
        try:
            # Iterate through each slide in the presentation
            for slide in self.presentation.slides:
                structure = {"placeholders": []}
                
                # Iterate through each shape in the slide
                for shape in slide.shapes:
                    # Safely check if the shape is a placeholder
                    is_placeholder = False
                    try:
                        if hasattr(shape, 'is_placeholder') and hasattr(shape, 'placeholder_format'):
                            is_placeholder = shape.is_placeholder
                    except Exception:
                        # If accessing placeholder properties fails, it's not a placeholder
                        is_placeholder = False
                    
                    if is_placeholder:
                        try:
                            # Create a dictionary to store placeholder details
                            placeholder = {
                                "type": shape.placeholder_format.type if hasattr(shape.placeholder_format, 'type') else 0,
                                "idx": shape.placeholder_format.idx if hasattr(shape.placeholder_format, 'idx') else 0,
                                "has_text_frame": shape.has_text_frame if hasattr(shape, 'has_text_frame') else False,
                                "name": shape.name if hasattr(shape, 'name') else "Unknown",
                            }
                            
                            # Check if the placeholder contains a table
                            has_table = False
                            if hasattr(shape, 'has_table'):
                                has_table = shape.has_table
                            
                            if has_table:
                                placeholder["has_table"] = True
                                # Extract the table structure as a list of rows with cell texts
                                try:
                                    placeholder["table_structure"] = [
                                        [cell.text for cell in row.cells] for row in shape.table.rows
                                    ]
                                except Exception:
                                    placeholder["table_structure"] = [[""]]
                            else:
                                placeholder["has_table"] = False
                            
                            # Check if the placeholder is an image placeholder
                            is_image = False
                            try:
                                if hasattr(shape, 'placeholder_format') and hasattr(shape.placeholder_format, 'type'):
                                    is_image = (shape.placeholder_format.type == 18)  # 18 is typically for images
                            except Exception:
                                is_image = False
                                
                            placeholder["has_image"] = is_image
                            placeholder["image_description"] = ""
        
                            # Append the placeholder information to the structure
                            structure["placeholders"].append(placeholder)
                        except Exception as e:
                            print(f"Warning: Error analyzing placeholder in template: {str(e)}")
                
                # Append the slide structure to the list of all slide structures
                slide_structures.append(structure)
        except Exception as e:
            print(f"Warning: Error analyzing PowerPoint template: {str(e)}")
            return []  # Return empty structure on error
            
        return slide_structures
    
    def generate_slide_from_text(self, title, content, layout_index=1):
        """
        Generate a new slide with title and content
        
        Args:
            title (str): Slide title
            content (str): Slide content
            layout_index (int, optional): Index of the slide layout to use
        
        Returns:
            The newly created slide
        """
        if layout_index >= len(self.presentation.slide_layouts):
            layout_index = 1  # Default to title and content if specified layout doesn't exist
            
        slide_layout = self.presentation.slide_layouts[layout_index]
        slide = self.presentation.slides.add_slide(slide_layout)
        
        # Set the title
        if hasattr(slide, 'shapes') and hasattr(slide.shapes, 'title') and slide.shapes.title:
            slide.shapes.title.text = title
        
        # Set the content
        for shape in slide.placeholders:
            if shape.placeholder_format.type == 1:  # 1 is title
                shape.text = title
            elif shape.placeholder_format.type == 7:  # 7 is body/content
                shape.text = content
        
        return slide

    def add_image_slide(self, title, image_path, layout_index=8):
        """
        Add a slide with an image
        
        Args:
            title (str): Slide title
            image_path (str): Path to the image file
            layout_index (int, optional): Index of the slide layout to use
        
        Returns:
            The newly created slide
        """
        if layout_index >= len(self.presentation.slide_layouts):
            layout_index = 8  # Default to picture with caption layout
            if layout_index >= len(self.presentation.slide_layouts):
                layout_index = 1  # Fallback to title and content
        
        slide_layout = self.presentation.slide_layouts[layout_index]
        slide = self.presentation.slides.add_slide(slide_layout)
        
        # Set the title
        if hasattr(slide, 'shapes') and hasattr(slide.shapes, 'title') and slide.shapes.title:
            slide.shapes.title.text = title
        
        # Add the image
        for shape in slide.placeholders:
            if shape.placeholder_format.type == 18:  # 18 is for pictures
                shape.insert_picture(image_path)
                break
        else:  # If no placeholder for image is found
            slide.shapes.add_picture(image_path, Inches(1.5), Inches(2), height=Inches(5))
        
        return slide
    
    def add_chart_slide(self, title, chart_data, chart_type='bar', layout_index=1):
        """
        Add a slide with a chart
        
        Args:
            title (str): Slide title
            chart_data (dict): Chart data in the format {categories: [], series: [{name: str, values: []}]}
            chart_type (str): Type of chart (bar, column, line, pie)
            layout_index (int): Index of the slide layout to use
        
        Returns:
            The newly created slide
        """
        slide_layout = self.presentation.slide_layouts[layout_index]
        slide = self.presentation.slides.add_slide(slide_layout)
        
        # Set the title
        if hasattr(slide, 'shapes') and hasattr(slide.shapes, 'title') and slide.shapes.title:
            slide.shapes.title.text = title
        
        # Create a temporary chart image using matplotlib
        chart_img_path = self._create_chart_image(title, chart_data, chart_type)
        
        # Add the chart image to the slide
        for shape in slide.placeholders:
            if shape.placeholder_format.type == 7:  # Content placeholder
                left = shape.left
                top = shape.top
                width = shape.width
                height = shape.height
                slide.shapes.add_picture(chart_img_path, left, top, width, height)
                break
        else:  # If no placeholder is found
            slide.shapes.add_picture(chart_img_path, Inches(1), Inches(2), width=Inches(8), height=Inches(4.5))
        
        # Clean up temporary file
        if os.path.exists(chart_img_path):
            os.remove(chart_img_path)
        
        return slide
    
    def _create_chart_image(self, title, chart_data, chart_type):
        """
        Create a chart image using matplotlib
        
        Args:
            title (str): Chart title
            chart_data (dict): Chart data
            chart_type (str): Type of chart
        
        Returns:
            str: Path to the temporary chart image file
        """
        plt.figure(figsize=(10, 6))
        
        # Extract categories and series data
        categories = chart_data.get('categories', [])
        series_list = chart_data.get('series', [])
        
        if chart_type.lower() == 'pie' and len(series_list) > 0:
            values = series_list[0].get('values', [])
            plt.pie(values, labels=categories, autopct='%1.1f%%', startangle=90)
            plt.axis('equal')  # Equal aspect ratio ensures that pie is drawn as a circle
            
        elif chart_type.lower() in ['bar', 'column']:
            x = np.arange(len(categories))
            width = 0.8 / len(series_list) if len(series_list) > 0 else 0.8
            
            for i, series in enumerate(series_list):
                values = series.get('values', [])
                if chart_type.lower() == 'bar':
                    plt.barh(x - width/2 + i*width, values, width, label=series.get('name', ''))
                else:  # column
                    plt.bar(x - width/2 + i*width, values, width, label=series.get('name', ''))
            
            plt.xticks(x, categories, rotation=45 if chart_type.lower() == 'column' else 0)
            plt.legend()
            
        elif chart_type.lower() == 'line':
            for series in series_list:
                values = series.get('values', [])
                plt.plot(categories, values, marker='o', label=series.get('name', ''))
            
            plt.xticks(rotation=45)
            plt.legend()
        
        plt.title(title)
        plt.tight_layout()
        
        # Save chart to a temporary file
        temp_file = tempfile.NamedTemporaryFile(delete=False, suffix='.png')
        plt.savefig(temp_file.name)
        plt.close()
        
        return temp_file.name
    
    def generate_from_data(self, data_path, output_path, title="Data Analysis Presentation", include_charts=True):
        """
        Generate a PowerPoint presentation from a data file (CSV, Excel)
        
        Args:
            data_path (str): Path to the data file
            output_path (str): Path where to save the output PowerPoint file
            title (str): Title of the presentation
            include_charts (bool): Whether to include charts visualizing the data
        """
        # Load the data
        file_ext = os.path.splitext(data_path)[1].lower()
        if file_ext == '.csv':
            df = pd.read_csv(data_path)
        elif file_ext in ['.xlsx', '.xls']:
            df = pd.read_excel(data_path)
        else:
            raise ValueError(f"Unsupported file format: {file_ext}")
        
        # Create title slide
        slide_layout = self.presentation.slide_layouts[0]  # Title slide layout
        slide = self.presentation.slides.add_slide(slide_layout)
        slide.shapes.title.text = title
        if len(slide.placeholders) > 1:
            # Typically the second placeholder is for subtitle
            subtitle = slide.placeholders[1]
            subtitle.text = f"Generated from {os.path.basename(data_path)}"
        
        # Add summary slide
        summary_text = f"""Dataset Summary:
        - Number of rows: {len(df)}
        - Number of columns: {len(df.columns)}
        - Columns: {', '.join(df.columns.tolist())}
        """
        self.generate_slide_from_text("Dataset Overview", summary_text)
        
        # Add data sample slide
        sample_df = df.head(10).copy()
        # Convert the sample df to string for display
        sample_table = sample_df.to_string(index=False)
        self.generate_slide_from_text("Data Sample", sample_table)
        
        if include_charts:
            # Generate charts for numeric columns
            numeric_cols = df.select_dtypes(include=[np.number]).columns.tolist()
            if numeric_cols:
                for col in numeric_cols[:3]:  # Limit to first 3 numeric columns
                    # Create bar chart
                    if df[col].nunique() < 15:  # Only for columns with reasonable number of unique values
                        value_counts = df[col].value_counts().nlargest(10)
                        chart_data = {
                            'categories': value_counts.index.astype(str).tolist(),
                            'series': [{'name': col, 'values': value_counts.values.tolist()}]
                        }
                        self.add_chart_slide(f"{col} Distribution", chart_data, 'column')
                
                # Create summary statistics slide
                stats_df = df[numeric_cols].describe().T
                stats_text = stats_df.to_string()
                self.generate_slide_from_text("Numeric Columns Statistics", stats_text)
            
            # Generate pie charts for categorical columns
            categorical_cols = df.select_dtypes(include=['object']).columns.tolist()
            if categorical_cols:
                for col in categorical_cols[:3]:  # Limit to first 3 categorical columns
                    if df[col].nunique() < 10:  # Only for columns with reasonable number of categories
                        value_counts = df[col].value_counts().nlargest(5)
                        chart_data = {
                            'categories': value_counts.index.tolist(),
                            'series': [{'name': col, 'values': value_counts.values.tolist()}]
                        }
                        self.add_chart_slide(f"{col} Distribution", chart_data, 'pie')
        
        # Save the presentation
        self.save(output_path)
    
    def generate_from_template(self, template_path, content_pdf_path, output_path, ai_content=None):
        """
        Generate a PowerPoint presentation using a PowerPoint template
        and filling it with content from a PDF using RAG
        
        Args:
            template_path (str): Path to the PowerPoint template (.pptx)
            content_pdf_path (str): Path to the PDF with content to extract using RAG
            output_path (str): Path where to save the output PowerPoint file
            ai_content (dict, optional): Pre-generated AI content to include
        """
        try:
            # Load the template if not already loaded
            if template_path and template_path.lower().endswith('.pptx') and \
               (not hasattr(self, 'presentation') or self.presentation is None):
                self.presentation = Presentation(template_path)
                self.slide_structures = self._analyze_template()
            
            # Process the content PDF with RAG
            if content_pdf_path and os.path.exists(content_pdf_path):
                chunks_added = self.rag_processor.process_pdf(content_pdf_path)
                print(f"Processed PDF and added {chunks_added} chunks to the vector database.")
            
            # Create title slide safely
            try:
                # Get first slide layout (usually title slide)
                if len(self.presentation.slide_layouts) > 0:
                    slide_layout = self.presentation.slide_layouts[0]
                    slide = self.presentation.slides.add_slide(slide_layout)
                    
                    # Get title from AI content or generate one
                    if ai_content and 'title' in ai_content:
                        title = ai_content['title']
                    else:
                        title = self.rag_processor.generate_content_with_rag("Generate a concise and engaging title for this presentation")
                        if not title or title.startswith("Error"):
                            title = os.path.splitext(os.path.basename(content_pdf_path))[0] if content_pdf_path else "AI Generated Presentation"
                    
                    # Try to set the title - with error handling
                    try:
                        # First try to access the title placeholder
                        if hasattr(slide.shapes, 'title') and slide.shapes.title is not None:
                            slide.shapes.title.text = title
                        else:
                            # Look for a title placeholder manually
                            title_set = False
                            for shape in slide.placeholders:
                                try:
                                    if shape.placeholder_format.type == 1:  # 1 is title
                                        shape.text = title
                                        title_set = True
                                        break
                                except:
                                    pass  # Not a valid placeholder
                            
                            # If no title placeholder found, add a text box
                            if not title_set:
                                tb = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(1.5))
                                tb.text_frame.text = title
                        
                        # Try to add a subtitle if placeholders exist
                        try:
                            if len(slide.placeholders) > 1:
                                subtitle = None
                                for shape in slide.placeholders:
                                    try:
                                        if shape.placeholder_format.type == 2:  # 2 is subtitle
                                            subtitle = shape
                                            break
                                    except:
                                        pass  # Not a valid placeholder
                                
                                if subtitle:
                                    subtitle.text = "AI-Generated Presentation"
                        except:
                            # Subtitle setting failed, but that's not critical
                            pass
                    except Exception as e:
                        print(f"Warning: Could not set title: {str(e)}")
                        # Add a textbox as a fallback
                        tb = slide.shapes.add_textbox(Inches(1), Inches(1), Inches(8), Inches(1.5))
                        tb.text_frame.text = title
            except Exception as e:
                print(f"Warning: Could not create title slide: {str(e)}")
            
            # Define common sections for presentations if we don't have a structured template
            sections = [
                "Introduction", "Overview", "Key Points", "Analysis",
                "Results", "Conclusions", "Recommendations", "Next Steps"
            ]
            
            # If we have slide structures from a template, try to use those instead
            if self.slide_structures:
                # Extract section names from slide structures (placeholder names)
                template_sections = []
                for slide_struct in self.slide_structures:
                    for placeholder in slide_struct.get("placeholders", []):
                        if placeholder.get("has_text_frame", False) and "title" in placeholder.get("name", "").lower():
                            section_name = placeholder.get("name", "")
                            if section_name and section_name not in ["Title 1", "Title"]:
                                template_sections.append(section_name)
                
                if template_sections:
                    sections = template_sections
            
            # Create slides for each section
            for section_name in sections:
                try:
                    # Generate content for this section using RAG
                    if ai_content and section_name.lower() in ai_content:
                        section_content = ai_content[section_name.lower()]
                    else:
                        # Generate content using RAG
                        query = f"Create presentation content for the '{section_name}' section based on the provided document."
                        section_content = self.rag_processor.generate_content_with_rag(query)
                    
                    # Create slide for this section if we have content
                    if section_content:
                        self.generate_slide_from_text(section_name, section_content)
                except Exception as e:
                    print(f"Warning: Error creating section '{section_name}': {str(e)}")
        except Exception as e:
            raise Exception(f"Error generating presentation from template: {str(e)}")
        
        # Save the presentation
        self.save(output_path)
    
    def _extract_text_from_pdf(self, pdf_path):
        """
        Extract text content from a PDF file
        
        Args:
            pdf_path (str): Path to the PDF file
        
        Returns:
            str: Extracted text content
        """
        with open(pdf_path, 'rb') as file:
            reader = PdfReader(file)
            text = ''
            for page in reader.pages:
                text += page.extract_text() + '\n\n'
            return text
    
    def _extract_sections_from_text(self, text):
        """
        Extract sections from text based on headings/formatting
        
        Args:
            text (str): Text to extract sections from
        
        Returns:
            dict: Dictionary of section name to placeholder content
        """
        # Simple regex-based section extraction
        # Look for lines that might be headings (all caps, followed by newline, etc.)
        section_pattern = r'([A-Z][A-Za-z\s]+:)|(^[A-Z][A-Z\s]+$)'
        sections = {}
        current_section = None
        current_content = []
        
        for line in text.split('\n'):
            line = line.strip()
            if not line:
                continue
                
            # Check if this line is a section heading
            matches = re.findall(section_pattern, line)
            if matches and any(match[0] or match[1] for match in matches):
                # Save previous section if exists
                if current_section:
                    sections[current_section] = '\n'.join(current_content)
                
                # Start new section
                current_section = line.strip(':').strip()
                current_content = []
            elif current_section:
                current_content.append(line)
        
        # Save the last section
        if current_section and current_content:
            sections[current_section] = '\n'.join(current_content)
        
        # If no sections found but there's content, use first line as title
        if not sections and text.strip():
            lines = [line for line in text.split('\n') if line.strip()]
            if lines:
                sections['title'] = lines[0]
                if len(lines) > 1:
                    sections['Content'] = '\n'.join(lines[1:])
        
        return sections
    
    def _find_section_content(self, content_text, section_name):
        """
        Try to find content for a specific section in the extracted text
        
        Args:
            content_text (str): Text content to search in
            section_name (str): Name of the section to find
        
        Returns:
            str: Content for the specified section, or empty string if not found
        """
        # Simple approach: look for section_name followed by text
        pattern = f"({section_name}:|{section_name.upper()}:|{section_name.title()}:)\s*([\s\S]+?)\n\n"
        matches = re.search(pattern, content_text)
        
        if matches and matches.group(2):
            return matches.group(2).strip()
            
        # If not found, try to match based on keywords
        keywords = section_name.lower().split()
        paragraphs = content_text.split('\n\n')
        
        for paragraph in paragraphs:
            if all(keyword in paragraph.lower() for keyword in keywords):
                return paragraph
        
        return ""
    
    def populate_from_ai(self, ai_content, output_path):
        """
        Populate a presentation with AI-generated content
        
        Args:
            ai_content (dict): Dictionary of section names to AI-generated content
            output_path (str): Path where to save the output PowerPoint file
        """
        # Create title slide
        slide_layout = self.presentation.slide_layouts[0]  # Title slide layout
        slide = self.presentation.slides.add_slide(slide_layout)
        
        slide.shapes.title.text = ai_content.get('title', 'AI Generated Presentation')
        if len(slide.placeholders) > 1 and 'subtitle' in ai_content:
            subtitle = slide.placeholders[1]
            subtitle.text = ai_content.get('subtitle', '')
        
        # Create content slides based on AI content
        for section_name, content in ai_content.items():
            if section_name in ['title', 'subtitle']:
                continue  # Already handled in title slide
            
            # If content is a string, create a simple slide
            if isinstance(content, str):
                self.generate_slide_from_text(section_name, content)
            
            # If content is a dictionary with special instructions
            elif isinstance(content, dict):
                if content.get('type') == 'chart':
                    self.add_chart_slide(section_name, content.get('data', {}), content.get('chart_type', 'bar'))
                elif content.get('type') == 'image' and 'path' in content:
                    self.add_image_slide(section_name, content['path'])
                else:
                    # Default to text slide with json content
                    self.generate_slide_from_text(section_name, json.dumps(content, indent=2))
        
        # Save the presentation
        self.save(output_path)
    
    def save(self, output_path):
        """
        Save the presentation to a file
        
        Args:
            output_path (str): Path where to save the presentation
        """
        # Ensure directory exists
        os.makedirs(os.path.dirname(os.path.abspath(output_path)), exist_ok=True)
        self.presentation.save(output_path)
