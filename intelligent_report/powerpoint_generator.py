# PowerPoint generation utilities
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.shapes import MSO_SHAPE
from pptx.dml.color import RGBColor
from typing import List, Dict, Any, Optional, Union


class PowerPointGenerator:
    """
    Handles creation of PowerPoint presentations from report data
    """
    
    def __init__(self, title: str = "Generated Report"):
        self.title = title
        self.subtitle = "Generated Report"
        self.author = ""
        self.company = ""
        self.prs = Presentation()
        
        # Add title slide
        self._add_title_slide()
    
    def set_metadata(self, subtitle: str = "", author: str = "", company: str = "") -> None:
        """Set metadata for the presentation"""
        self.subtitle = subtitle if subtitle else "Generated Report"
        self.author = author
        self.company = company
        
        # Update the title slide
        if len(self.prs.slides) > 0:
            self._update_title_slide()
    
    def _add_title_slide(self) -> None:
        """Add the title slide to the presentation"""
        title_slide_layout = self.prs.slide_layouts[0]  # Title slide layout
        slide = self.prs.slides.add_slide(title_slide_layout)
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        
        title.text = self.title
        subtitle.text = self.subtitle
        
        # Add company and author info if available
        if self.author or self.company:
            txBox = slide.shapes.add_textbox(Inches(0.5), Inches(5), Inches(9), Inches(0.75))
            tf = txBox.text_frame
            
            if self.company:
                p = tf.add_paragraph()
                p.text = self.company
                p.font.size = Pt(14)
            
            if self.author:
                p = tf.add_paragraph()
                p.text = f"Prepared by: {self.author}"
                p.font.size = Pt(12)
    
    def _update_title_slide(self) -> None:
        """Update the title slide with current metadata"""
        if len(self.prs.slides) == 0:
            return
            
        # Get the title slide (first slide)
        slide = self.prs.slides[0]
        
        # Update title and subtitle
        for shape in slide.shapes:
            if shape.has_text_frame:
                if shape.text.strip() == "Generated Report" or "Click to edit Master subtitle style" in shape.text:
                    shape.text = self.subtitle
                # Title is usually already set correctly, but update it just in case
                elif "Click to edit Master title style" in shape.text:
                    shape.text = self.title
        
        # Remove any existing company/author text boxes
        shapes_to_remove = []
        for i, shape in enumerate(slide.shapes):
            if shape.has_text_frame and ("Prepared by:" in shape.text or self.company in shape.text):
                shapes_to_remove.append(i)
        
        # We can't remove shapes while iterating, so we gather indices and delete after
        # Note: PowerPoint doesn't allow direct removal of shapes by index, so we skip this step
        
        # Add new company and author info
        if self.author or self.company:
            txBox = slide.shapes.add_textbox(Inches(0.5), Inches(5), Inches(9), Inches(0.75))
            tf = txBox.text_frame
            
            if self.company:
                p = tf.add_paragraph()
                p.text = self.company
                p.font.size = Pt(14)
            
            if self.author:
                p = tf.add_paragraph()
                p.text = f"Prepared by: {self.author}"
                p.font.size = Pt(12)
    
    def add_section(self, section_name: str, section_data: Dict) -> None:
        """Add a section to the presentation"""
        # Create section slide with heading and content
        content_slide_layout = self.prs.slide_layouts[1]  # Title and content layout
        slide = self.prs.slides.add_slide(content_slide_layout)
        
        # Set the title
        title = slide.shapes.title
        title.text = section_name
        
        # Add content
        content = slide.placeholders[1]
        
        # Format the content based on the section type
        if section_data["type"] == "user":
            content.text = section_data["content"]
        elif section_data["type"] == "ai":
            # Limit AI-generated content to fit in a slide
            text = section_data["content"]
            # Split into paragraphs and limit to first 3-4 paragraphs or ~150 words
            paragraphs = text.split('\n\n')
            content.text = '\n'.join(paragraphs[:min(3, len(paragraphs))])
            
            # If content is too long, add a "continued..." slide
            if len(paragraphs) > 3:
                slide = self.prs.slides.add_slide(content_slide_layout)
                title = slide.shapes.title
                title.text = f"{section_name} (Continued)"
                content = slide.placeholders[1]
                content.text = '\n'.join(paragraphs[3:min(6, len(paragraphs))])
        elif section_data["type"] == "rag":
            # For RAG content, include the main content and sources
            text_parts = section_data["content"].split("Sources:")
            if len(text_parts) > 1:
                # Main content
                main_text = text_parts[0].strip()
                # Sources
                sources = text_parts[1].strip()
                
                # Add main content
                paragraphs = main_text.split('\n\n')
                content.text = '\n'.join(paragraphs[:min(3, len(paragraphs))])
                
                # Add sources in a bullet list
                txBox = slide.shapes.add_textbox(Inches(0.5), Inches(5), Inches(9), Inches(1.5))
                tf = txBox.text_frame
                tf.text = "Sources:"
                
                # Add sources as bullet points
                for source_line in sources.split('\n'):
                    p = tf.add_paragraph()
                    p.text = source_line.strip()
                    p.level = 1  # Set bullet level
            else:
                # If no sources, just add the content
                content.text = section_data["content"]
    
    def add_visualization(self, viz_data: Dict) -> None:
        """Add a visualization to the presentation"""
        # Create slide with just a title
        title_only_slide_layout = self.prs.slide_layouts[5]  # Title only layout
        slide = self.prs.slides.add_slide(title_only_slide_layout)
        
        # Set the title
        title = slide.shapes.title
        title.text = viz_data['title']
        
        # Add description if available
        if viz_data.get('description'):
            txBox = slide.shapes.add_textbox(Inches(0.5), Inches(1.2), Inches(9), Inches(0.5))
            tf = txBox.text_frame
            tf.text = viz_data['description']
        
        # Add the visualization image
        if os.path.exists(viz_data['filename']):
            left = Inches(1.5)
            top = Inches(2)
            width = Inches(7)  # Adjust based on image aspect ratio
            slide.shapes.add_picture(viz_data['filename'], left, top, width=width)
    
    def generate(self, output_path: str) -> str:
        """Generate and save the PowerPoint presentation"""
        # Add a closing slide
        self._add_closing_slide()
        
        # Ensure directory exists
        os.makedirs(os.path.dirname(output_path), exist_ok=True)
        
        # Save the presentation
        self.prs.save(output_path)
        print(f"PowerPoint presentation generated successfully: {output_path}")
        return output_path
    
    def _add_closing_slide(self) -> None:
        """Add a closing slide to the presentation"""
        layout = self.prs.slide_layouts[6]  # Blank slide layout
        slide = self.prs.slides.add_slide(layout)
        
        # Add title text box
        title = slide.shapes.add_textbox(Inches(1), Inches(2), Inches(8), Inches(1.5))
        title_frame = title.text_frame
        title_frame.text = "Thank You"
        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(44)
        title_para.font.bold = True
        title_para.alignment = 1  # Center alignment
