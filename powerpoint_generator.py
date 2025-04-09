import json
import os
import re
import requests
import matplotlib.pyplot as plt
import numpy as np
import pdfplumber
from pptx import Presentation
from pptx.util import Inches
from ollama import Client
from chromadb import PersistentClient
from chromadb.utils import embedding_functions
import tempfile
import logging
import shutil

class PowerPointGenerator:
    def __init__(self, template_path, data_path):
        self.template_path = os.path.join(os.path.dirname(__file__), template_path)
        self.data_path = data_path
        self.presentation = Presentation(self.template_path)
        self.client = Client()
        self.model = "mistral"
        self.chroma_client = PersistentClient(path="./chroma_db")
        self._setup_graph_config()
        
    def _setup_graph_config(self):
        """Configure graph settings."""
        self.graph_config = {
            "incident_response": {
                "title": "Incident Response Time Analysis",
                "xlabel": "Quarter",
                "ylabel": "Time (minutes/hours)",
                "metrics": [
                    {"name": "avg_times", "label": "Average Response Time", "color": "blue", "marker": "o"},
                    {"name": "median_times", "label": "Median Response Time", "color": "orange", "marker": "s"},
                    {"name": "resolution_times", "label": "Resolution Time", "color": "green", "marker": "^"}
                ]
            },
            "security_events": {
                "severity": {
                    "title": "Severity Distribution",
                    "labels": ["Critical", "High", "Medium", "Low"],
                    "colors": ["red", "orange", "yellow", "green"]
                },
                "types": {
                    "title": "Event Type Distribution",
                    "labels": ["Phishing", "Malware", "DDoS", "Unauthorized Access"]
                }
            },
            "patch_management": {
                "title": "Patch Management Status",
                "ylabel": "Percentage",
                "metrics": [
                    {"name": "deployed", "color": "green", "label": "Deployed"},
                    {"name": "pending", "color": "red", "label": "Pending"}
                ]
            },
            "user_behavior": {
                "auth": {
                    "title": "Authentication Success Rate",
                    "labels": ["Success", "Failed"],
                    "colors": ["green", "red"]
                },
                "mfa": {
                    "title": "MFA Usage Rate",
                    "labels": ["With MFA", "Without MFA"],
                    "colors": ["blue", "orange"]
                }
            },
            "vulnerabilities": {
                "title": "Vulnerability Distribution by System",
                "systems": ["Web Servers", "Databases", "Network Devices", "Workstations"],
                "severity_levels": ["Critical", "High", "Medium", "Low"]
            }
        }

    def setup_chroma_db(self):
        """Set up ChromaDB collection and ingest PDF data."""
        # Create or get collection
        collection = self.chroma_client.get_or_create_collection(
            name="cybersecurity_data",
            metadata={"hnsw:space": "cosine"}
        )

        # Ingest PDF data
        with pdfplumber.open(self.data_path) as pdf:
            text = ""
            for page in pdf.pages:
                text += page.extract_text()
            
        # Split text into chunks
        chunks = [text[i:i+1000] for i in range(0, len(text), 1000)]
        
        # Add to collection
        collection.add(
            documents=chunks,
            ids=[f"chunk_{i}" for i in range(len(chunks))]
        )
        
        return collection

    def analyze_template_structure(self):
        """Analyze the PowerPoint template structure."""
        slide_structures = []
        for slide in self.presentation.slides:
            structure = {"placeholders": []}
            for shape in slide.shapes:
                if shape.is_placeholder:
                    placeholder = {
                        "type": shape.placeholder_format.type,
                        "idx": shape.placeholder_format.idx,
                        "has_text_frame": shape.has_text_frame,
                        "name": shape.name,
                        "has_table": shape.has_table if hasattr(shape, 'has_table') else False,
                        "has_image": shape.placeholder_format.type == 18
                    }
                    structure["placeholders"].append(placeholder)
            slide_structures.append(structure)
        return slide_structures

    def extract_data_from_pdf(self):
        """Extract structured data from the PDF using LLM."""
        data = {
            "incident_response": {},
            "security_events": {},
            "patch_management": {},
            "user_behavior": {},
            "vulnerabilities": {}
        }
        
        try:
            # Extract text from PDF
            with pdfplumber.open(self.data_path) as pdf:
                text = "\n".join(page.extract_text() for page in pdf.pages)
                
            # Use LLM to categorize and extract data
            prompt = f"""Extract cybersecurity metrics from the following text. 
            Categorize the data into these sections:
            1. Incident Response (response times, resolution times)
            2. Security Events (event types, severity levels)
            3. Patch Management (patch status, deployment times)
            4. User Behavior (authentication rates, MFA usage)
            5. Vulnerabilities (system vulnerabilities, severity levels)
            
            Return the data in JSON format:
            ```json
            {{
                "incident_response": {{
                    "avg_times": [...],
                    "median_times": [...],
                    "resolution_times": [...]
                }},
                "security_events": {{
                    "severity": [...],
                    "counts": [...],
                    "event_types": [...],
                    "type_counts": [...]
                }},
                "patch_management": {{
                    "severity_levels": [...],
                    "deployed": [...],
                    "pending": [...]
                }},
                "user_behavior": {{
                    "auth_success": 0.0,
                    "mfa_usage": 0.0,
                    "failed_attempts": 0.0
                }},
                "vulnerabilities": {{
                    "systems": [...],
                    "severity_levels": [...],
                    "counts": [...]
                }}
            }}
            ```
            
            Text to analyze:
            {text}
            """

            response = self.client.chat(
                model=self.model,
                messages=[{"role": "user", "content": prompt}],
                stream=False
            )
            
            # Parse the response
            content = response["message"]["content"]
            try:
                extracted_data = json.loads(content)
                data.update(extracted_data)
            except json.JSONDecodeError:
                print("Warning: Could not parse JSON from LLM response. Using default data.")
                # Fallback to default data
                data["incident_response"] = {
                    "avg_times": [60, 50, 45, 40],
                    "median_times": [55, 45, 40, 35],
                    "resolution_times": [6, 5, 4, 5.5]
                }
                data["security_events"] = {
                    "severity": ["Critical", "High", "Medium", "Low"],
                    "counts": [120, 250, 500, 380],
                    "event_types": ["Phishing", "Malware", "DDoS", "Unauthorized Access"],
                    "type_counts": [150, 200, 100, 120]
                }
                data["patch_management"] = {
                    "severity_levels": ["Critical", "High", "Medium"],
                    "deployed": [100, 95, 85],
                    "pending": [0, 5, 15]
                }
                data["user_behavior"] = {
                    "auth_success": 99.2,
                    "mfa_usage": 85,
                    "failed_attempts": 0.8
                }
                data["vulnerabilities"] = {
                    "systems": ["Web Servers", "Databases", "Network Devices", "Workstations"],
                    "severity_levels": ["Critical", "High", "Medium", "Low"],
                    "counts": [
                        [5, 10, 15, 20],
                        [10, 15, 20, 25],
                        [15, 20, 25, 30],
                        [20, 25, 30, 35]
                    ]
                }

        except Exception as e:
            print(f"Error extracting data: {str(e)}")
            # Fallback to default data
            data["incident_response"] = {
                "avg_times": [60, 50, 45, 40],
                "median_times": [55, 45, 40, 35],
                "resolution_times": [6, 5, 4, 5.5]
            }
            data["security_events"] = {
                "severity": ["Critical", "High", "Medium", "Low"],
                "counts": [120, 250, 500, 380],
                "event_types": ["Phishing", "Malware", "DDoS", "Unauthorized Access"],
                "type_counts": [150, 200, 100, 120]
            }
            data["patch_management"] = {
                "severity_levels": ["Critical", "High", "Medium"],
                "deployed": [100, 95, 85],
                "pending": [0, 5, 15]
            }
            data["user_behavior"] = {
                "auth_success": 99.2,
                "mfa_usage": 85,
                "failed_attempts": 0.8
            }
            data["vulnerabilities"] = {
                "systems": ["Web Servers", "Databases", "Network Devices", "Workstations"],
                "severity_levels": ["Critical", "High", "Medium", "Low"],
                "counts": [
                    [5, 10, 15, 20],
                    [10, 15, 20, 25],
                    [15, 20, 25, 30],
                    [20, 25, 30, 35]
                ]
            }

        return data

    def generate_content(self, topic, structure):
        """Generate slide content using LLM based on the extracted data."""
        # Get the extracted data
        data = self.extract_data_from_pdf()
        
        # Define the slide structure
        slides = []
        
        # Generate content for each metric category
        for category in ["incident_response", "security_events", "patch_management", "user_behavior", "vulnerabilities"]:
            # Create a prompt for each category
            prompt = f"""Generate PowerPoint slide content for {category} metrics.
            
            Data context:
            {json.dumps(data[category], indent=2)}
            
            Create 3 slides with these components:
            1. Title that clearly describes the metric
            2. Content that explains the data and its significance
            3. Description of the graph to be generated
            
            Return in this format:
            ```json
            [
                {{
                    "placeholders": [
                        {{
                            "name": "Title 1",
                            "text": "Your Title Here"
                        }},
                        {{
                            "name": "Content Placeholder 2",
                            "text": "Your content here"
                        }},
                        {{
                            "name": "Image Placeholder 3",
                            "image_description": "Your graph description"
                        }}
                    ]
                }}
            ]
            ```
            """

            response = self.client.chat(
                model=self.model,
                messages=[{"role": "user", "content": prompt}],
                stream=False
            )
            
            try:
                content = response["message"]["content"]
                category_slides = json.loads(content)
                slides.extend(category_slides)
            except json.JSONDecodeError:
                print(f"Warning: Could not parse JSON for {category}. Using default content.")
                # Add default slides for this category
                for i in range(3):
                    slides.append({
                        "placeholders": [
                            {
                                "name": "Title 1",
                                "text": f"{category.replace('_', ' ').title()} - Slide {i+1}"
                            },
                            {
                                "name": "Content Placeholder 2",
                                "text": f"Analysis of {category.replace('_', ' ')} metrics."
                            },
                            {
                                "name": "Image Placeholder 3",
                                "image_description": f"{category.replace('_', ' ')} visualization"
                            }
                        ]
                    })

        return slides, data

    def _generate_incident_response_graph(self, ax, data):
        """Generate incident response time graph."""
        config = self.graph_config["incident_response"]
        labels = ['Q1', 'Q2', 'Q3', 'Q4']
        
        for metric in config["metrics"]:
            values = data["incident_response"][metric["name"]]
            ax.plot(labels, values, 
                   marker=metric["marker"], 
                   label=metric["label"], 
                   color=metric["color"])
        
        ax.set_title(config["title"])
        ax.set_xlabel(config["xlabel"])
        ax.set_ylabel(config["ylabel"])
        ax.legend()
        
        # Add value labels
        for i in range(len(labels)):
            for metric in config["metrics"]:
                values = data["incident_response"][metric["name"]]
                y = values[i]
                ax.text(i, y + 2, f'{y}m', 
                       ha='center', va='bottom', 
                       color=metric["color"])

    def _generate_security_events_graph(self, fig, data):
        """Generate security events distribution graph."""
        config = self.graph_config["security_events"]
        ax1, ax2 = fig.subplots(1, 2)
        
        # Severity distribution
        counts = data["security_events"]["counts"]
        ax1.pie(counts, 
               labels=config["severity"]["labels"], 
               colors=config["severity"]["colors"], 
               autopct='%1.1f%%',
               startangle=90)
        ax1.set_title(config["severity"]["title"])
        
        # Event types
        event_counts = data["security_events"]["type_counts"]
        ax2.bar(config["types"]["labels"], event_counts, color='blue')
        ax2.set_title(config["types"]["title"])
        ax2.set_ylabel('Number of Events')
        
        # Add value labels
        for i, v in enumerate(event_counts):
            ax2.text(i, v + 5, str(v), ha='center', va='bottom')

    def _generate_patch_management_graph(self, ax, data):
        """Generate patch management status graph."""
        config = self.graph_config["patch_management"]
        levels = data["patch_management"]["severity_levels"]
        
        for metric in config["metrics"]:
            values = data["patch_management"][metric["name"]]
            if metric["name"] == "deployed":
                ax.bar(levels, values, color=metric["color"], label=metric["label"])
            else:
                deployed = data["patch_management"]["deployed"]
                ax.bar(levels, values, 
                       bottom=deployed, 
                       color=metric["color"], 
                       label=metric["label"])
        
        ax.set_title(config["title"])
        ax.set_ylabel(config["ylabel"])
        ax.legend()
        
        # Add value labels
        for i in range(len(levels)):
            deployed = data["patch_management"]["deployed"][i]
            pending = data["patch_management"]["pending"][i]
            ax.text(i, deployed/2, f'{deployed}%', 
                   ha='center', va='center', color='white')
            ax.text(i, deployed + pending/2, f'{pending}%', 
                   ha='center', va='center', color='white')

    def _generate_user_behavior_graph(self, fig, data):
        """Generate user behavior analysis graph."""
        config = self.graph_config["user_behavior"]
        ax1, ax2 = fig.subplots(1, 2)
        
        # Authentication success rate
        auth_counts = [
            data["user_behavior"]["auth_success"],
            data["user_behavior"]["failed_attempts"]
        ]
        ax1.pie(auth_counts, 
               labels=config["auth"]["labels"], 
               autopct='%1.1f%%',
               colors=config["auth"]["colors"],
               startangle=90)
        ax1.set_title(config["auth"]["title"])
        
        # MFA usage
        mfa_counts = [
            data["user_behavior"]["mfa_usage"],
            100 - data["user_behavior"]["mfa_usage"]
        ]
        ax2.pie(mfa_counts, 
               labels=config["mfa"]["labels"], 
               autopct='%1.1f%%',
               colors=config["mfa"]["colors"],
               startangle=90)
        ax2.set_title(config["mfa"]["title"])

    def _generate_vulnerability_graph(self, ax, data):
        """Generate vulnerability assessment graph."""
        config = self.graph_config["vulnerabilities"]
        counts = data["vulnerabilities"]["counts"]
        
        im = ax.imshow(counts, cmap='Reds')
        ax.set_xticks(np.arange(len(config["systems"])))
        ax.set_yticks(np.arange(len(config["severity_levels"])))
        ax.set_xticklabels(config["systems"])
        ax.set_yticklabels(config["severity_levels"])
        
        plt.setp(ax.get_xticklabels(), rotation=45, ha="right", rotation_mode="anchor")
        
        # Add value labels
        for i in range(len(config["severity_levels"])):
            for j in range(len(config["systems"])):
                ax.text(j, i, counts[i, j],
                       ha="center", va="center", color="white")
        
        ax.set_title(config["title"])

    def generate_graph(self, query, data):
        """Generate a graph based on the query and data context."""
        try:
            # Create a temporary directory for graphs
            temp_dir = tempfile.mkdtemp(prefix="pptx_graphs_")
            graph_path = os.path.join(temp_dir, "temp_graph.png")
            
            # Create figure with appropriate size
            fig, ax = plt.subplots(figsize=(12, 8))
            
            # Generate graph based on query
            if "incident response" in query.lower():
                self._generate_incident_response_graph(ax, data)
            elif "security events" in query.lower():
                self._generate_security_events_graph(fig, data)
            elif "patch management" in query.lower():
                self._generate_patch_management_graph(ax, data)
            elif "user behavior" in query.lower():
                self._generate_user_behavior_graph(fig, data)
            elif "vulnerability" in query.lower():
                self._generate_vulnerability_graph(ax, data)
            else:
                # Default bar chart
                labels = ['Metric 1', 'Metric 2', 'Metric 3', 'Metric 4']
                values = [np.random.randint(50, 100) for _ in range(4)]
                ax.bar(labels, values, color='blue')
                ax.set_title('General Metric Distribution')
                ax.set_ylabel('Value')
                
                # Add value labels
                for i, v in enumerate(values):
                    ax.text(i, v + 1, str(v), ha='center', va='bottom')
            
            # Save and close the figure
            plt.savefig(graph_path, bbox_inches='tight', dpi=300)
            plt.close()
            
            return graph_path
            
        except Exception as e:
            logging.error(f"Error generating graph: {str(e)}")
            raise

    def populate_slide(self, slide, slide_data, data):
        """Populate a single slide with generated content."""
        placeholders = slide_data.get("placeholders", [])
        
        # Populate text placeholders
        for placeholder in placeholders:
            for shape in slide.shapes:
                if shape.is_placeholder and placeholder.get("name") in shape.name:
                    if placeholder.get("text") and shape.has_text_frame:
                        shape.text = placeholder["text"]

        # Handle image placeholders
        for placeholder in placeholders:
            if placeholder.get("image_description"):
                # Generate graph based on description and data
                graph_path = self.generate_graph(placeholder["image_description"], data)
                
                # Find the image placeholder
                for shape in slide.shapes:
                    if shape.is_placeholder and shape.placeholder_format.type == 18:
                        # Remove existing image placeholder
                        slide.shapes._spTree.remove(shape._element)
                        
                        # Add new image
                        try:
                            img = slide.shapes.add_picture(
                                graph_path,
                                shape.left,
                                shape.top,
                                shape.width,
                                shape.height
                            )
                            
                            # Ensure the image is properly positioned
                            img.left = shape.left
                            img.top = shape.top
                            img.width = shape.width
                            img.height = shape.height
                            
                            # Clean up the temporary image file
                            if os.path.exists(graph_path):
                                os.remove(graph_path)
                            
                            break
                        except Exception as e:
                            print(f"Error adding image: {str(e)}")
                            # Add a text box explaining the error
                            txBox = slide.shapes.add_textbox(
                                shape.left,
                                shape.top,
                                shape.width,
                                shape.height
                            )
                            tf = txBox.text_frame
                            tf.text = f"Error displaying graph: {str(e)}"
                            break

        # Handle tables
        for placeholder in placeholders:
            table_data = placeholder.get("table_structure")
            if table_data:
                for shape in slide.shapes:
                    if shape.has_table:
                        table = shape.table
                        for i, row_data in enumerate(table_data):
                            for j, cell_text in enumerate(row_data):
                                table.cell(i, j).text = cell_text
                        break

    def generate_presentation(self, topic, output_path):
        """Generate a complete PowerPoint presentation."""
        try:
            # Analyze template structure
            structure = self.analyze_template_structure()
            
            # Generate content and get data
            slides, data = self.generate_content(topic, structure)
            
            if not slides:
                print("Failed to generate valid content")
                return False
            
            # Create a temporary directory
            temp_dir = tempfile.mkdtemp()
            temp_pptx = os.path.join(temp_dir, "temp_output.pptx")
            
            # Create a new presentation
            presentation = Presentation()
            
            # Add title slide
            title_slide_layout = presentation.slide_layouts[0]
            slide = presentation.slides.add_slide(title_slide_layout)
            title = slide.shapes.title
            subtitle = slide.placeholders[1]
            
            title.text = topic
            subtitle.text = "Cybersecurity Metrics Analysis"
            
            # Add content slides
            for i, slide_data in enumerate(slides):
                # Use appropriate layout based on content
                if i == 0:  # Title slide
                    layout = presentation.slide_layouts[0]
                elif i == 1:  # Introduction slide
                    layout = presentation.slide_layouts[1]
                else:  # Content slides
                    layout = presentation.slide_layouts[5]  # Blank layout
                
                slide = presentation.slides.add_slide(layout)
                self.populate_slide(slide, slide_data, data)
            
            # Save to temporary location first
            presentation.save(temp_pptx)
            
            # Move to final output location
            if output_path:
                shutil.move(temp_pptx, output_path)
            else:
                output_path = "output.pptx"
                shutil.move(temp_pptx, output_path)
            
            print("Presentation generated successfully!")
            return True
            
        except Exception as e:
            print(f"Error generating presentation: {str(e)}")
            return False

if __name__ == "__main__":
    # Example usage
    generator = PowerPointGenerator(
        "templates/template.pptx",
        "sample_data.pdf"
    )
    generator.collection = generator.setup_chroma_db()
    success = generator.generate_presentation(
        "Cybersecurity Metrics Analysis",
        "output.pptx"
    )
    
    if success:
        print("Presentation generated successfully!")
    else:
        print("Failed to generate presentation.")
